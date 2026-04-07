import io
import os
import time
from collections import defaultdict
from contextlib import asynccontextmanager

from fastapi import FastAPI, HTTPException, Depends, Query, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.security import OAuth2PasswordRequestForm
from pydantic import BaseModel
from sqlalchemy.orm import Session
from openai import OpenAI
from dotenv import load_dotenv

from app import models, schemas
from app.dependencies import get_db
from app.auth import (
    hash_password,
    verify_password,
    create_access_token,
    get_current_user,
)

load_dotenv()

API_KEY = os.getenv("GROQ_API_KEY")
model_context = {}

# ── Rate Limiting ────────────────────────────────────────────────────────────
ADMIN_EMAILS = {"yuvraj@nandini.com", "abhinavnair@gmail.com"}
RATE_LIMIT_REQUESTS = 30        # max requests
RATE_LIMIT_WINDOW_SEC = 3600    # per hour

_request_log: dict[str, list[float]] = defaultdict(list)


def _apply_rate_limit(email: str) -> None:
    """Raise 429 if a non-admin user exceeds the hourly AI request quota."""
    if email in ADMIN_EMAILS:
        return
    now = time.time()
    window_start = now - RATE_LIMIT_WINDOW_SEC
    log = _request_log[email]
    # Prune stale entries
    _request_log[email] = [t for t in log if t > window_start]
    if len(_request_log[email]) >= RATE_LIMIT_REQUESTS:
        raise HTTPException(
            status_code=429,
            detail=f"Rate limit exceeded: {RATE_LIMIT_REQUESTS} AI requests per hour. Try again later.",
        )
    _request_log[email].append(now)


# ── Lifespan (Groq Client Init) ──────────────────────────────────────────────
@asynccontextmanager
async def lifespan(app: FastAPI):
    print("🔌 Initializing NeuroNotes Pro...")
    if not API_KEY:
        print("❌ ERROR: 'GROQ_API_KEY' not found in .env!")
    else:
        try:
            client = OpenAI(
                api_key=API_KEY,
                base_url="https://api.groq.com/openai/v1",
            )
            model_context["client"] = client
            print("🚀 Groq AI Client successfully initialized!")
        except Exception as e:
            print(f"❌ CRITICAL ERROR: {e}")
    yield
    model_context.clear()


app = FastAPI(lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://127.0.0.1:8000", "http://localhost:8000"],
    allow_methods=["*"],
    allow_headers=["*"],
)


# ── Request Schemas ──────────────────────────────────────────────────────────
class GenerateRequest(BaseModel):
    system_prompt: str | None = None
    prompt: str
    max_tokens: int = 4000
    temperature: float = 0.7


class ChatRequest(BaseModel):
    message: str
    context: str = ""
    temperature: float = 0.3


# ── AUTH ENDPOINTS ───────────────────────────────────────────────────────────

@app.post("/register", response_model=schemas.UserResponse)
def register(user: schemas.UserCreate, db: Session = Depends(get_db)):
    existing_user = (
        db.query(models.User)
        .filter(models.User.email == user.email)
        .first()
    )

    if existing_user:
        raise HTTPException(status_code=400, detail="Email already registered")

    new_user = models.User(
        email=user.email,
        hashed_password=hash_password(user.password),
    )

    db.add(new_user)
    db.commit()
    db.refresh(new_user)

    return new_user


@app.post("/login")
def login(
    form_data: OAuth2PasswordRequestForm = Depends(),
    db: Session = Depends(get_db),
):
    db_user = (
        db.query(models.User)
        .filter(models.User.email == form_data.username)
        .first()
    )

    if not db_user or not verify_password(
        form_data.password,
        db_user.hashed_password,
    ):
        raise HTTPException(status_code=401, detail="Invalid credentials")

    access_token = create_access_token(
        data={"sub": db_user.email}
    )

    return {
        "access_token": access_token,
        "token_type": "bearer",
    }


@app.get("/me")
def read_current_user(
    current_user: models.User = Depends(get_current_user),
):
    return {
        "id": current_user.id,
        "email": current_user.email,
        "is_active": current_user.is_active,
    }


@app.post("/change-password")
def change_password(
    payload: schemas.ChangePasswordRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    if not verify_password(payload.old_password, current_user.hashed_password):
        raise HTTPException(
            status_code=400, detail="Old password is incorrect")

    current_user.hashed_password = hash_password(payload.new_password)

    db.commit()
    db.refresh(current_user)

    return {"message": "Password updated successfully"}


# ── NOTES CRUD ───────────────────────────────────────────────────────────────

@app.post("/notes", response_model=schemas.NoteResponse)
def create_note(
    note: schemas.NoteCreate,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    new_note = models.Note(
        title=note.title,
        content=note.content,
        owner_id=current_user.id,
    )

    db.add(new_note)
    db.commit()
    db.refresh(new_note)

    return new_note


@app.get("/notes", response_model=list[schemas.NoteResponse])
def get_notes(
    saved: bool | None = Query(default=None),
    search: str | None = Query(default=None),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    query = (
        db.query(models.Note)
        .filter(models.Note.owner_id == current_user.id)
    )

    if saved is True:
        query = query.filter(models.Note.is_bookmarked == True)

    if search:
        term = f"%{search}%"
        query = query.filter(
            models.Note.title.ilike(term) | models.Note.content.ilike(term)
        )

    notes = query.order_by(models.Note.created_at.desc()).all()

    return notes


@app.patch("/notes/{note_id}", response_model=schemas.NoteResponse)
def update_note(
    note_id: int,
    updates: schemas.NoteUpdate,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    note = (
        db.query(models.Note)
        .filter(
            models.Note.id == note_id,
            models.Note.owner_id == current_user.id,
        )
        .first()
    )

    if not note:
        raise HTTPException(status_code=404, detail="Note not found")

    if updates.title is not None:
        note.title = updates.title
    if updates.content is not None:
        note.content = updates.content

    db.commit()
    db.refresh(note)

    return note


@app.delete("/notes/{note_id}")
def delete_note(
    note_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    note = (
        db.query(models.Note)
        .filter(
            models.Note.id == note_id,
            models.Note.owner_id == current_user.id,
        )
        .first()
    )

    if not note:
        raise HTTPException(status_code=404, detail="Note not found")

    db.delete(note)
    db.commit()

    return {"message": "Note deleted successfully"}


@app.patch("/notes/{note_id}/bookmark")
def toggle_bookmark(
    note_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    note = (
        db.query(models.Note)
        .filter(
            models.Note.id == note_id,
            models.Note.owner_id == current_user.id,
        )
        .first()
    )

    if not note:
        raise HTTPException(status_code=404, detail="Note not found")

    note.is_bookmarked = not note.is_bookmarked

    db.commit()
    db.refresh(note)

    return note


# ── AI ENDPOINTS ─────────────────────────────────────────────────────────────

@app.post("/generate")
async def generate_text(
    request: GenerateRequest,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    _apply_rate_limit(current_user.email)

    client = model_context.get("client")
    if not client:
        raise HTTPException(
            status_code=503, detail="AI Client not initialized.")

    try:
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {
                    "role": "system",
                    "content": request.system_prompt or "You are NeuroNotes Pro."
                },
                {"role": "user", "content": request.prompt},
            ],
            temperature=request.temperature,
            max_tokens=request.max_tokens,
        )

        generated_text = completion.choices[0].message.content

        new_note = models.Note(
            title=(request.prompt.strip().splitlines()[0])[:50] if request.prompt else "Untitled",
            content=generated_text,
            owner_id=current_user.id,
            is_bookmarked=False,
        )

        db.add(new_note)
        db.commit()
        db.refresh(new_note)

        return {
            "response": generated_text,
            "note_id": new_note.id,
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Groq API Error: {str(e)}")


@app.post("/chat")
async def chat(
    request: ChatRequest,
    current_user: models.User = Depends(get_current_user),
):
    """Contextual chat — calls the AI but does NOT save a note."""
    _apply_rate_limit(current_user.email)

    client = model_context.get("client")
    if not client:
        raise HTTPException(
            status_code=503, detail="AI Client not initialized.")

    try:
        system_prompt = (
            "You are a helpful study assistant. "
            "Use the provided context to answer the user's question. "
            "Keep your answer concise, conversational, and format it in markdown."
        )
        prompt = f"Context:\n{request.context}\n\nUser Question:\n{request.message}"

        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt},
            ],
            temperature=request.temperature,
            max_tokens=1500,
        )

        return {"response": completion.choices[0].message.content}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Groq API Error: {str(e)}")


# ── FILE UPLOAD ───────────────────────────────────────────────────────────────

MAX_UPLOAD_BYTES = 10 * 1024 * 1024  # 10 MB


@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    current_user: models.User = Depends(get_current_user),
):
    """Extract plain text from an uploaded PDF or PPTX file."""
    content = await file.read()

    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail="File too large. Maximum size is 10 MB.",
        )

    filename = (file.filename or "").lower()

    if filename.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n\n".join(p.strip() for p in pages if p.strip())
            if not text:
                raise HTTPException(
                    status_code=422,
                    detail="Could not extract text from this PDF. It may be image-only or encrypted.",
                )
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"PDF read error: {str(e)}")

    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if slide_texts:
                    parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
            if not text:
                raise HTTPException(
                    status_code=422,
                    detail="No text found in this presentation.",
                )
        except HTTPException:
            raise
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"PPTX read error: {str(e)}")

    elif filename.endswith(".ppt"):
        raise HTTPException(
            status_code=415,
            detail="Old .ppt format is not supported. Please save as .pptx and re-upload.",
        )

    elif filename.endswith((".png", ".jpg", ".jpeg")):
        raise HTTPException(
            status_code=415,
            detail="Image OCR is not supported. Please copy the text from the image manually.",
        )

    else:
        raise HTTPException(
            status_code=415,
            detail="Unsupported file type. Supported formats: PDF, PPTX.",
        )

    return {"text": text, "filename": file.filename}


# ── FLASHCARDS ────────────────────────────────────────────────────────────────

@app.post("/flashcards", response_model=schemas.FlashcardDeckResponse)
def save_flashcard_deck(
    deck: schemas.FlashcardDeckCreate,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    new_deck = models.FlashcardDeck(
        topic=deck.topic,
        difficulty=deck.difficulty,
        count=len(deck.cards),
        cards=deck.cards,
        owner_id=current_user.id,
    )
    db.add(new_deck)
    db.commit()
    db.refresh(new_deck)
    return schemas.FlashcardDeckResponse(
        id=new_deck.id,
        topic=new_deck.topic,
        difficulty=new_deck.difficulty,
        count=new_deck.count,
        cards=new_deck.cards,
        saved_at=new_deck.saved_at.isoformat(),
    )


@app.get("/flashcards", response_model=list[schemas.FlashcardDeckResponse])
def get_flashcard_decks(
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    decks = (
        db.query(models.FlashcardDeck)
        .filter(models.FlashcardDeck.owner_id == current_user.id)
        .order_by(models.FlashcardDeck.saved_at.desc())
        .all()
    )
    return [
        schemas.FlashcardDeckResponse(
            id=d.id,
            topic=d.topic,
            difficulty=d.difficulty,
            count=d.count,
            cards=d.cards,
            saved_at=d.saved_at.isoformat(),
        )
        for d in decks
    ]


@app.delete("/flashcards/{deck_id}")
def delete_flashcard_deck(
    deck_id: int,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    deck = (
        db.query(models.FlashcardDeck)
        .filter(
            models.FlashcardDeck.id == deck_id,
            models.FlashcardDeck.owner_id == current_user.id,
        )
        .first()
    )
    if not deck:
        raise HTTPException(status_code=404, detail="Deck not found")
    db.delete(deck)
    db.commit()
    return {"message": "Deck deleted"}


# ── STATIC FILES ──────────────────────────────────────────────────────────────
app.mount("/", StaticFiles(directory=".", html=True), name="static")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(
        "main:app",
        host="127.0.0.1",
        port=8000,
        reload=True,
        reload_excludes=[".venv/*", "venv/*", "__pycache__/*", ".git/*"],
    )
